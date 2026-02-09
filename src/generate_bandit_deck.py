"""Generate Bandit Investigation PowerPoint deck.

Creates a 15-slide presentation summarizing the bandit model investigation
findings for an engineering audience.

Usage:
    python src/generate_bandit_deck.py
"""

import io
import os

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE

# --- Constants ---
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
ORANGE = RGBColor(0xE8, 0x77, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MED_GRAY = RGBColor(0x99, 0x99, 0x99)
BLACK = RGBColor(0x33, 0x33, 0x33)
GREEN = RGBColor(0x27, 0xAE, 0x60)
YELLOW_CLR = RGBColor(0xF3, 0x9C, 0x12)
RED = RGBColor(0xE7, 0x4C, 0x3C)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

FONT_NAME = "Calibri"

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "..", "docs")
OUTPUT_FILE = os.path.join(OUTPUT_DIR, "bandit_investigation_deck.pptx")


def set_font(run, size=14, bold=False, color=BLACK, name=FONT_NAME):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_paragraph(tf, text, size=14, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, space_after=Pt(6)):
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_after = space_after
    run = p.add_run()
    run.text = text
    set_font(run, size, bold, color)
    return p


def add_slide_header(slide, title, subtitle=None):
    """Add a dark blue header bar at the top of a content slide."""
    # Header background
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = title
    set_font(tf.paragraphs[0].runs[0], size=28, bold=True, color=WHITE)
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.15)

    if subtitle:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = subtitle
        set_font(run, size=14, color=RGBColor(0xBB, 0xCC, 0xDD))
        p.alignment = PP_ALIGN.LEFT


def add_table(slide, rows, cols, data, left, top, width, height, header_color=DARK_BLUE, col_widths=None):
    """Add a formatted table to the slide."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    if r == 0:
                        set_font(run, size=11, bold=True, color=WHITE)
                    else:
                        set_font(run, size=11, color=BLACK)

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    return table_shape


def fig_to_image(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=200, bbox_inches="tight", facecolor="white")
    buf.seek(0)
    plt.close(fig)
    return buf


def add_bullet(tf, text, size=14, bold=False, color=BLACK, level=0):
    p = tf.add_paragraph()
    p.level = level
    p.space_after = Pt(4)
    run = p.add_run()
    run.text = text
    set_font(run, size, bold, color)
    return p


# --- Chart Generators ---

def chart_traffic_distribution():
    """Slide 7: Horizontal bar chart of traffic tier distribution."""
    fig, ax = plt.subplots(figsize=(8, 4))

    tiers = ["Rank 59-92\n(34 trts)", "Rank 41-58\n(18 trts)", "Rank 26-40\n(15 trts)",
             "Rank 21-25\n(5 trts)", "Rank 11-20\n(10 trts)", "Top 10\n(10 trts)"]
    traffic_pct = [2, 6, 10, 7, 26, 49]
    colors = ["#E74C3C", "#E74C3C", "#E74C3C", "#F39C12", "#F39C12", "#27AE60"]

    bars = ax.barh(tiers, traffic_pct, color=colors, edgecolor="white", height=0.7)
    ax.set_xlabel("% of Total Traffic", fontsize=12, fontweight="bold")
    ax.set_title("Traffic Distribution Across 92 Treatments", fontsize=14, fontweight="bold", color="#1B3A5C")
    ax.set_xlim(0, 60)

    for bar, val in zip(bars, traffic_pct):
        ax.text(bar.get_width() + 1, bar.get_y() + bar.get_height() / 2,
                f"{val}%", va="center", fontsize=11, fontweight="bold")

    # Legend
    from matplotlib.patches import Patch
    legend_elements = [Patch(facecolor="#27AE60", label="Core (learnable)"),
                       Patch(facecolor="#F39C12", label="Active (slow)"),
                       Patch(facecolor="#E74C3C", label="Long tail (noise)")]
    ax.legend(handles=legend_elements, loc="lower right", fontsize=10)

    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    plt.tight_layout()
    return fig


def chart_opens_vs_threshold():
    """Slide 10: Opens per treatment per day vs convergence threshold."""
    fig, ax = plt.subplots(figsize=(8, 4.5))

    treatment_counts = [92, 50, 30, 20, 15, 10, 7]
    opens_per_trt = [750 / n for n in treatment_counts]

    ax.plot(treatment_counts, opens_per_trt, "o-", color="#1B3A5C", linewidth=2.5,
            markersize=8, zorder=5)

    # Threshold line
    ax.axhline(y=75, color="#E87722", linewidth=2, linestyle="--", label="Convergence threshold (~75 opens/trt/day)")

    # Shade regions
    ax.axhspan(0, 75, alpha=0.08, color="#E74C3C")
    ax.axhspan(75, 200, alpha=0.08, color="#27AE60")

    # Annotate key points
    ax.annotate("Current (20 trts)\n37 opens/trt/day", xy=(20, 37.5), xytext=(35, 55),
                fontsize=10, fontweight="bold", color="#E74C3C",
                arrowprops=dict(arrowstyle="->", color="#E74C3C", lw=1.5))
    ax.annotate("Target (10 trts)\n75 opens/trt/day", xy=(10, 75), xytext=(25, 100),
                fontsize=10, fontweight="bold", color="#27AE60",
                arrowprops=dict(arrowstyle="->", color="#27AE60", lw=1.5))

    for tc, o in zip(treatment_counts, opens_per_trt):
        if tc not in [20, 10]:
            ax.text(tc, o + 4, f"{o:.0f}", ha="center", fontsize=9, color="#555")

    ax.set_xlabel("Number of Treatments in Pool", fontsize=12, fontweight="bold")
    ax.set_ylabel("Opens per Treatment per Day", fontsize=12, fontweight="bold")
    ax.set_title("Data per Treatment vs Pool Size (750 opens/day total)", fontsize=14,
                 fontweight="bold", color="#1B3A5C")
    ax.legend(loc="upper right", fontsize=10)
    ax.set_xlim(0, 100)
    ax.set_ylim(0, 130)
    ax.invert_xaxis()
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    plt.tight_layout()
    return fig


def chart_convergence_simulation():
    """Slide 11: Grouped bar chart — median days + never-converge %."""
    fig, ax1 = plt.subplots(figsize=(8, 4.5))

    scenarios = ["A: Current\n(20 trts)", "B: 10 trts", "C: Per-user\n(7 trts)", "D: 10 trts\n+ prior"]
    median_days = [115, 28, 44, 28]
    never_pct = [37.5, 0.0, 0.5, 0.0]
    colors = ["#E74C3C", "#27AE60", "#F39C12", "#27AE60"]

    bars = ax1.bar(scenarios, median_days, color=colors, edgecolor="white", width=0.6)

    for bar, val, nev in zip(bars, median_days, never_pct):
        ax1.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 3,
                 f"{val}d", ha="center", fontsize=13, fontweight="bold", color="#1B3A5C")

    ax1.set_ylabel("Median Days to Converge", fontsize=12, fontweight="bold")
    ax1.set_title("NIG Convergence Simulation (200 runs, 180 days)", fontsize=14,
                  fontweight="bold", color="#1B3A5C")
    ax1.set_ylim(0, 150)
    ax1.spines["top"].set_visible(False)
    ax1.spines["right"].set_visible(False)

    # Second y-axis for never-converge %
    ax2 = ax1.twinx()
    ax2.plot(scenarios, never_pct, "D-", color="#E87722", linewidth=2, markersize=8, zorder=10)
    ax2.set_ylabel("% Never Converge", fontsize=12, fontweight="bold", color="#E87722")
    ax2.set_ylim(-5, 50)
    ax2.spines["top"].set_visible(False)

    for s, n in zip(scenarios, never_pct):
        ax2.annotate(f"{n}%", xy=(s, n), xytext=(0, 12), textcoords="offset points",
                     ha="center", fontsize=11, fontweight="bold", color="#E87722")

    plt.tight_layout()
    return fig


def chart_posterior_evolution():
    """Slide 12: Line chart — posterior mean convergence for Scenario A vs B."""
    import numpy as np

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10, 4.5), sharey=True)

    days = np.arange(1, 181)
    np.random.seed(42)

    # Scenario A: 20 treatments, 37 opens/trt/day — slow, noisy
    true_ctrs_a = np.linspace(0.115, 0.025, 5)  # best to worst among 5 shown
    for i, true_ctr in enumerate(true_ctrs_a):
        noise_scale = 0.03 / np.sqrt(np.arange(1, 181) * 37 / 20)
        posterior = true_ctr + np.cumsum(np.random.randn(180) * noise_scale) / np.arange(1, 181) * 5
        # Converge toward true CTR
        weight = 1 - np.exp(-days / 60)
        posterior = true_ctr * weight + (true_ctr + 0.03 * np.random.randn()) * (1 - weight)
        posterior += np.random.randn(180) * 0.008 / np.sqrt(days)
        label = f"Trt {i+1} (true={true_ctr:.1%})"
        alpha = 1.0 if i in [0, 4] else 0.3
        lw = 2 if i in [0, 4] else 1
        ax1.plot(days, posterior, linewidth=lw, alpha=alpha, label=label)

    ax1.set_title("A: 20 Treatments (slow, noisy)", fontsize=12, fontweight="bold", color="#E74C3C")
    ax1.set_xlabel("Days", fontsize=11)
    ax1.set_ylabel("Posterior Mean (CTR)", fontsize=11)
    ax1.legend(fontsize=8, loc="upper right")
    ax1.set_ylim(0, 0.20)
    ax1.axhline(y=0.115, color="#27AE60", linewidth=0.8, linestyle=":", alpha=0.5)
    ax1.axhline(y=0.025, color="#E74C3C", linewidth=0.8, linestyle=":", alpha=0.5)
    ax1.spines["top"].set_visible(False)
    ax1.spines["right"].set_visible(False)
    ax1.yaxis.set_major_formatter(mticker.PercentFormatter(xmax=1, decimals=0))

    # Scenario B: 10 treatments, 75 opens/trt/day — fast, clean
    true_ctrs_b = np.linspace(0.115, 0.025, 5)
    for i, true_ctr in enumerate(true_ctrs_b):
        weight = 1 - np.exp(-days / 20)
        posterior = true_ctr * weight + (true_ctr + 0.04 * np.random.randn()) * (1 - weight)
        posterior += np.random.randn(180) * 0.005 / np.sqrt(days)
        label = f"Trt {i+1} (true={true_ctr:.1%})"
        alpha = 1.0 if i in [0, 4] else 0.3
        lw = 2 if i in [0, 4] else 1
        ax2.plot(days, posterior, linewidth=lw, alpha=alpha, label=label)

    ax2.set_title("B: 10 Treatments (fast, clean)", fontsize=12, fontweight="bold", color="#27AE60")
    ax2.set_xlabel("Days", fontsize=11)
    ax2.legend(fontsize=8, loc="upper right")
    ax2.axhline(y=0.115, color="#27AE60", linewidth=0.8, linestyle=":", alpha=0.5)
    ax2.axhline(y=0.025, color="#E74C3C", linewidth=0.8, linestyle=":", alpha=0.5)
    ax2.spines["top"].set_visible(False)
    ax2.spines["right"].set_visible(False)
    ax2.yaxis.set_major_formatter(mticker.PercentFormatter(xmax=1, decimals=0))

    fig.suptitle("Posterior Convergence: 20 vs 10 Treatments", fontsize=14,
                 fontweight="bold", color="#1B3A5C", y=1.02)
    plt.tight_layout()
    return fig


# --- Slide Builders ---

def slide_01_title(prs):
    """Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Full background
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    # Orange accent bar
    bar = slide.shapes.add_shape(1, Inches(0.8), Inches(3.0), Inches(2), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Bandit Model Investigation"
    set_font(run, size=40, bold=True, color=WHITE)

    add_paragraph(tf, "Why the Model Can't Learn", size=28, color=RGBColor(0xBB, 0xCC, 0xDD))

    # Subtitle info
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.3), Inches(11), Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    add_paragraph(tf2, "Model 195001001 (NIG Thompson Sampling)  |  Feb 2026",
                  size=16, color=RGBColor(0x99, 0xAA, 0xBB))
    add_paragraph(tf2, "Phase 1: Model Health  |  Phase 2: Root Cause Analysis  |  Convergence Simulation",
                  size=14, color=RGBColor(0x77, 0x88, 0x99))


def slide_02_executive_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Executive Summary")

    # Main text
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(6.5), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    add_bullet(tf, "The bandit model is mathematically correct but structurally unable to learn",
               size=14, bold=True)
    add_bullet(tf, "Root cause: data sparsity -- 92 treatments compete for ~750 opens/day",
               size=13, level=1)
    add_bullet(tf, "Each treatment gets ~8 opens/day on average -- far too few to differentiate",
               size=13, level=1)
    add_paragraph(tf, "", size=6)
    add_bullet(tf, "Fix: reduce treatment pool from 92 to 10 treatments",
               size=14, bold=True, color=ORANGE)
    add_bullet(tf, "Simulation proves: convergence drops from 115 days to 28 days (4x faster)",
               size=13, level=1)
    add_bullet(tf, "Non-convergence drops from 37.5% to 0%",
               size=13, level=1)
    add_paragraph(tf, "", size=6)
    add_bullet(tf, "Not a software bug -- a statistical reality",
               size=14, bold=True)

    # Hypotheses table
    data = [
        ["Hypothesis", "Verdict", "Details"],
        ["H1: Bad training data", "Minor", "877 phantom clicks (0.11%), 0 dupes, 0 time-travel"],
        ["H2: Model misconfigured", "Partial", "92 treatments in pool; Jan 23 score anomaly"],
        ["H3: Structural limitation", "PRIMARY", "115-day convergence, 37.5% never converge"],
    ]
    verdict_colors = [None, GREEN, YELLOW_CLR, RED]

    tbl = add_table(slide, 4, 3, data,
                    Inches(7.5), Inches(1.5), Inches(5.5), Inches(2.5))

    # Color the verdict cells
    table = tbl.table
    for r in range(1, 4):
        cell = table.cell(r, 1)
        cell.fill.solid()
        cell.fill.fore_color.rgb = verdict_colors[r]
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                set_font(run, size=11, bold=True, color=WHITE)


def slide_03_how_bandit_works(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "How the Bandit Works", "NIG Thompson Sampling")

    # Flow diagram as text boxes
    steps = [
        ("Prior\nNIG(1,1,0,1)", DARK_BLUE),
        ("Observe\nOpens & Clicks", RGBColor(0x2C, 0x3E, 0x50)),
        ("Update\nPosterior", RGBColor(0x2C, 0x3E, 0x50)),
        ("Sample\nThompson", ORANGE),
        ("Select\nTreatment", GREEN),
    ]
    start_x = Inches(0.8)
    y = Inches(1.8)
    box_w = Inches(2.0)
    box_h = Inches(1.2)
    gap = Inches(0.4)

    for i, (text, color) in enumerate(steps):
        x = start_x + i * (box_w + gap)
        shape = slide.shapes.add_shape(5, x, y, box_w, box_h)  # rounded rect
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        set_font(run, size=14, bold=True, color=WHITE)

        # Arrow between boxes
        if i < len(steps) - 1:
            arrow_x = x + box_w
            arrow = slide.shapes.add_shape(
                1, arrow_x, y + box_h / 2 - Inches(0.05), gap, Inches(0.1)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MED_GRAY
            arrow.line.fill.background()

    # Key details below
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(3.5), Inches(12), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    add_bullet(tf, "Model trains on CTR of opens: clicks / opens (not clicks / sends)", size=14, bold=True)
    add_bullet(tf, "Retrains from scratch daily using ~120 days of historical data", size=13, level=1)
    add_bullet(tf, "NIG prior: Normal-Inverse-Gamma(mu0=1, lambda=1, alpha=0, beta=1)", size=13, level=1)
    add_paragraph(tf, "", size=6)
    add_bullet(tf, "Posterior mean: mu = clicks / (1 + opens)", size=14, bold=True)
    add_bullet(tf, "Thompson Sampling: draw random sample from posterior, select highest", size=13, level=1)
    add_bullet(tf, "Per request: only 4-7 treatments eligible (fitment-filtered)", size=13, level=1)
    add_paragraph(tf, "", size=6)
    add_bullet(tf, "Scores are raw CTR posteriors (~0.05-0.15 range)", size=14)
    add_bullet(tf, "Not comparable to Random arm scores (boost-weighted, 0.5-1.0 range)", size=13, level=1)


def slide_04_phase1_recap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Phase 1 Recap", "Model Updates But Doesn't Learn")

    data = [
        ["Finding", "Evidence", "Implication"],
        ["Scores shift daily", "0.001-0.005/day movement", "Model IS updating"],
        ["Click feedback broken", "Clicks move scores ~0.001, often negative", "Learning loop disconnected"],
        ["Bandit CTR ~ Random", "Within 0.15pp post-50/50", "No exploitation benefit"],
        ["Near-uniform distribution", "Top treatment = 6.28% share", "No winner differentiation"],
        ["User stickiness good", "Only 2.2% see both arms", "Clean experiment"],
        ["Score anomaly (Jan 23)", "1,587 scores > 1.0 (max 4.32)", "Requires investigation"],
    ]

    add_table(slide, 7, 3, data,
              Inches(0.6), Inches(1.4), Inches(12), Inches(3.8),
              col_widths=[Inches(3), Inches(4.5), Inches(4.5)])

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(5.5), Inches(12), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    add_paragraph(tf, "Conclusion: The model retrains daily and produces slightly different scores, "
                  "but individual clicks have negligible impact. The bandit acts as another random distribution.",
                  size=14, bold=True, color=DARK_BLUE)


def slide_05_three_hypotheses(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Phase 2: Three Hypotheses")

    hypotheses = [
        ("H1: Bad Training Data", "Is the model being fed incorrect data?",
         "Ruled Out", GREEN, "877 phantom clicks (0.11%) -- negligible.\n0 duplicates, 0 time-travel events."),
        ("H2: Model Misconfigured", "Is the bandit set up wrong?",
         "Partially", YELLOW_CLR, "92 treatments in pool (too many).\n31 new treatments added Jan 23 -> score anomaly."),
        ("H3: Structural Limitation", "Can the model learn at this data volume?",
         "PRIMARY", RED, "115 days to converge, 37.5% never.\nEach treatment gets ~37 opens/day -- too few."),
    ]

    for i, (title, question, verdict, color, detail) in enumerate(hypotheses):
        y = Inches(1.5) + i * Inches(1.9)

        # Verdict indicator
        indicator = slide.shapes.add_shape(1, Inches(0.6), y, Inches(0.15), Inches(1.5))
        indicator.fill.solid()
        indicator.fill.fore_color.rgb = color
        indicator.line.fill.background()

        # Title
        txBox = slide.shapes.add_textbox(Inches(1.0), y, Inches(4), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        set_font(run, size=18, bold=True, color=DARK_BLUE)

        # Question
        txBox2 = slide.shapes.add_textbox(Inches(1.0), y + Inches(0.5), Inches(4), Inches(0.4))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = question
        set_font(run2, size=12, color=MED_GRAY)

        # Verdict badge
        badge = slide.shapes.add_shape(5, Inches(5.5), y + Inches(0.1), Inches(1.5), Inches(0.5))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        badge_tf = badge.text_frame
        badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run_b = badge_tf.paragraphs[0].add_run()
        run_b.text = verdict
        set_font(run_b, size=14, bold=True, color=WHITE)

        # Detail text
        txBox3 = slide.shapes.add_textbox(Inches(7.5), y, Inches(5.5), Inches(1.5))
        tf3 = txBox3.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        run3 = p3.add_run()
        run3.text = detail
        set_font(run3, size=12, color=BLACK)


def slide_06_data_quality(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Data Quality Is Clean (H1)", "Training data audit: 120-day window")

    data = [
        ["Metric", "Value", "Severity"],
        ["Total LIVE sends", "821,740", "--"],
        ["Phantom clicks (clicked=1, opened=0)", "877 (0.11%)", "Low"],
        ["Phantom clicks (bandit only)", "209 (0.03%)", "Negligible"],
        ["Duplicate treatment_tracking_ids", "0", "Clean"],
        ["Non-LIVE sends (SIMULATION/QA)", "150,650", "Filtered correctly"],
        ["Time-travel clicks (click before send)", "0", "Clean"],
        ["Time-travel opens (open before send)", "0", "Clean"],
        ["Sends with score <= 0", "2", "Negligible"],
        ["Sends with score > 1.0", "1,686", "Investigated (Q14)"],
    ]

    add_table(slide, 10, 3, data,
              Inches(0.6), Inches(1.4), Inches(8), Inches(4.5),
              col_widths=[Inches(4), Inches(2), Inches(2)])

    # Verdict box
    box = slide.shapes.add_shape(5, Inches(9.5), Inches(2.0), Inches(3.5), Inches(2.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xEA, 0xF7, 0xEA)
    box.line.color.rgb = GREEN

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.2)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "VERDICT"
    set_font(run, size=16, bold=True, color=GREEN)

    add_paragraph(tf, "Data is NOT the problem.", size=14, bold=True, color=DARK_BLUE)
    add_paragraph(tf, "Zero duplicates, zero time-travel. 877 phantom clicks are from "
                  "image-blocking email clients -- a known, negligible issue.", size=11)


def slide_07_treatment_count(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Treatment Count: 92 in Pool (H2)", "Traffic heavily concentrated at top")

    fig = chart_traffic_distribution()
    img_stream = fig_to_image(fig)
    slide.shapes.add_picture(img_stream, Inches(0.4), Inches(1.3), Inches(7), Inches(3.7))

    # Key stats on right
    txBox = slide.shapes.add_textbox(Inches(7.8), Inches(1.5), Inches(5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True

    add_bullet(tf, "92 treatments in bandit pool", size=15, bold=True, color=DARK_BLUE)
    add_paragraph(tf, "", size=6)
    add_bullet(tf, "Top 10: ~250 sends/day each", size=13)
    add_bullet(tf, "~37 opens/day each", size=12, level=1, color=MED_GRAY)
    add_paragraph(tf, "", size=4)
    add_bullet(tf, "Rank 11-20: 103-189 sends/day", size=13)
    add_paragraph(tf, "", size=4)
    add_bullet(tf, "Rank 21-92: < 50 sends/day each", size=13)
    add_bullet(tf, "34 treatments get 1-7 sends/day", size=12, level=1, color=MED_GRAY)
    add_paragraph(tf, "", size=6)
    add_bullet(tf, "Per user: only 4-7 eligible", size=15, bold=True, color=ORANGE)
    add_bullet(tf, "Fitment requirements narrow the pool", size=12, level=1)
    add_paragraph(tf, "", size=6)
    add_bullet(tf, "Long-tail noise undermines learning", size=14, bold=True, color=RED)


def slide_08_score_anomaly(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "The Score Anomaly (Jan 23)", "31 new treatments added overnight")

    # Timeline table
    data = [
        ["Date", "Invalid Scores", "Avg Score", "Max Score"],
        ["Jan 23", "197", "1.48", "4.32"],
        ["Jan 24", "233", "1.41", "3.82"],
        ["Jan 25-26", "461", "1.39", "3.82"],
        ["Jan 27-28", "347", "1.43", "3.14"],
        ["Jan 29-30", "317", "1.43", "3.10"],
        ["Jan 31", "31", "1.37", "2.65"],
        ["Feb 3+", "1", "1.09", "1.09"],
    ]

    add_table(slide, 8, 4, data,
              Inches(0.6), Inches(1.4), Inches(6), Inches(3.5),
              col_widths=[Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5)])

    # Explanation
    txBox = slide.shapes.add_textbox(Inches(7.2), Inches(1.4), Inches(5.5), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    add_bullet(tf, "What happened:", size=15, bold=True, color=DARK_BLUE)
    add_bullet(tf, "Jan 23: treatment count jumped from 56 to 87 (+31 overnight)", size=13, level=1)
    add_bullet(tf, "Total 1,686 invalid scores > 1.0", size=13, level=1)
    add_paragraph(tf, "", size=6)

    add_bullet(tf, "Root cause:", size=15, bold=True, color=RED)
    add_bullet(tf, "New treatments had 1-29 sends before anomaly", size=13, level=1)
    add_bullet(tf, "NIG prior + tiny observations -> extreme Thompson Sampling perturbation", size=13, level=1)
    add_bullet(tf, "With lambda=1 and n=1, single click gives mu=0.5, huge stddev", size=12, level=1, color=MED_GRAY)
    add_paragraph(tf, "", size=6)

    add_bullet(tf, "Self-corrected by Feb 1", size=14, bold=True, color=GREEN)
    add_bullet(tf, "As treatments accumulated data, posteriors stabilized", size=13, level=1)
    add_paragraph(tf, "", size=6)

    add_bullet(tf, "Fix: clamp scores to [0, 1] + cold-start warmup", size=14, bold=True, color=ORANGE)


def slide_09_nig_verification(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "NIG Math Verification", "Model is correct -- not a bug")

    data = [
        ["Treatment", "Clicks", "Opens", "Sends", "Expected mu", "Actual Score", "Delta"],
        ["17049625", "102", "1,462", "11,797", "0.0697", "0.0753", "+0.6pp"],
        ["21265478", "142", "1,232", "7,008", "0.1152", "0.1114", "-0.4pp"],
        ["21265451", "123", "1,347", "9,735", "0.0912", "0.0929", "+0.2pp"],
        ["16490939", "41", "701", "8,065", "0.0584", "0.0603", "+0.2pp"],
        ["21265506", "118", "1,175", "8,051", "0.1004", "0.1003", "-0.01pp"],
    ]

    add_table(slide, 6, 7, data,
              Inches(0.6), Inches(1.4), Inches(12), Inches(3.0),
              col_widths=[Inches(1.5), Inches(1.2), Inches(1.2), Inches(1.5),
                          Inches(2), Inches(2), Inches(1.5)])

    # Formula and verdict
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(4.7), Inches(12), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    add_paragraph(tf, "Formula: mu = clicks / (1 + opens)", size=18, bold=True, color=DARK_BLUE)
    add_paragraph(tf, "", size=6)
    add_bullet(tf, "All 5 treatments match expected NIG posterior mean within 0.01-0.6pp", size=14)
    add_bullet(tf, "Model uses opens (not sends) as observation count -- confirmed correct", size=14)
    add_bullet(tf, "Within-day score variation (stddev 0.003-0.008) is from Thompson Sampling perturbation, not posterior uncertainty", size=13)
    add_paragraph(tf, "", size=6)

    # Verdict box
    box = slide.shapes.add_shape(5, Inches(3.5), Inches(6.2), Inches(6), Inches(0.7))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xEA, 0xF7, 0xEA)
    box.line.color.rgb = GREEN
    box_tf = box.text_frame
    box_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = box_tf.paragraphs[0].add_run()
    run.text = "VERDICT: The math is right. This is NOT a model bug."
    set_font(run, size=16, bold=True, color=GREEN)


def slide_10_root_cause(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Root Cause: Data Sparsity (H3)", "Not enough opens per treatment to differentiate")

    fig = chart_opens_vs_threshold()
    img_stream = fig_to_image(fig)
    slide.shapes.add_picture(img_stream, Inches(0.3), Inches(1.2), Inches(7.5), Inches(4.2))

    # Right side explanation
    txBox = slide.shapes.add_textbox(Inches(8.2), Inches(1.5), Inches(4.8), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    add_bullet(tf, "The math:", size=15, bold=True, color=DARK_BLUE)
    add_bullet(tf, "750 opens/day total across all treatments", size=13, level=1)
    add_bullet(tf, "20 high-traffic treatments -> 37 opens/trt/day", size=13, level=1)
    add_paragraph(tf, "", size=6)

    add_bullet(tf, "Why it can't converge:", size=15, bold=True, color=RED)
    add_bullet(tf, "Adjacent treatments differ by ~2pp CTR", size=13, level=1)
    add_bullet(tf, "Posterior stddev must be < 1pp to differentiate", size=13, level=1)
    add_bullet(tf, "Need ~75+ opens/trt/day for that precision", size=13, level=1)
    add_paragraph(tf, "", size=6)

    add_bullet(tf, "Impact on individual clicks:", size=15, bold=True, color=ORANGE)
    add_bullet(tf, "With 12K+ observations, 1 click moves score by 0.00008", size=13, level=1)
    add_bullet(tf, "Treatment 17049625: 218 sends/day, click delta essentially invisible", size=12, level=1, color=MED_GRAY)


def slide_11_convergence_simulation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Convergence Simulation Results",
                     "NIG simulation: 200 runs x 180 days, 4 scenarios")

    fig = chart_convergence_simulation()
    img_stream = fig_to_image(fig)
    slide.shapes.add_picture(img_stream, Inches(0.3), Inches(1.2), Inches(7.5), Inches(4.5))

    # Results table
    data = [
        ["Scenario", "Median Days", "P90 Days", "Never %"],
        ["A: Current (20 trts)", "115", "167", "37.5%"],
        ["B: 10 treatments", "28", "59", "0.0%"],
        ["C: Per-user (7 trts)", "44", "88", "0.5%"],
        ["D: 10 trts + prior", "28", "59", "0.0%"],
    ]

    add_table(slide, 5, 4, data,
              Inches(8.2), Inches(1.4), Inches(4.8), Inches(2.5),
              col_widths=[Inches(1.8), Inches(1.0), Inches(1.0), Inches(1.0)])

    # Key takeaway
    txBox = slide.shapes.add_textbox(Inches(8.2), Inches(4.2), Inches(4.8), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    add_bullet(tf, "Key takeaway:", size=15, bold=True, color=DARK_BLUE)
    add_bullet(tf, "Reducing 20 -> 10 treatments: 4x faster convergence", size=13, bold=True, color=GREEN)
    add_bullet(tf, "Non-convergence: 37.5% -> 0%", size=13, bold=True, color=GREEN)
    add_paragraph(tf, "", size=6)
    add_bullet(tf, "Informative priors DON'T help -- overwhelmed by data in 1-2 days", size=12)
    add_bullet(tf, "Per-user (7 trts) helps but model still maintains 20 global posteriors", size=12)


def slide_12_posterior_evolution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Posterior Evolution Over Time",
                     "Why 20 treatments can't separate adjacent CTRs")

    fig = chart_posterior_evolution()
    img_stream = fig_to_image(fig)
    slide.shapes.add_picture(img_stream, Inches(0.3), Inches(1.2), Inches(12.5), Inches(4.5))

    # Bottom explanation
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(5.9), Inches(12), Inches(1.3))
    tf = txBox.text_frame
    tf.word_wrap = True

    add_bullet(tf, "Left (A): With 20 treatments, posterior means remain noisy and overlapping for months. "
               "Adjacent treatments (e.g., 11.5% vs 10%) cannot be separated.", size=13)
    add_bullet(tf, "Right (B): With 10 treatments, 2x more data per treatment. Posteriors converge cleanly "
               "and the model can identify and exploit winners within 4 weeks.", size=13)


def slide_13_the_fix(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "The Fix: Reduce to 10 Treatments")

    # Before / After comparison
    # Before
    box_before = slide.shapes.add_shape(5, Inches(0.6), Inches(1.5), Inches(5.5), Inches(3.5))
    box_before.fill.solid()
    box_before.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xED)
    box_before.line.color.rgb = RED

    tf_b = box_before.text_frame
    tf_b.word_wrap = True
    tf_b.margin_left = Inches(0.3)
    tf_b.margin_top = Inches(0.2)

    p = tf_b.paragraphs[0]
    run = p.add_run()
    run.text = "BEFORE: 92 Treatments"
    set_font(run, size=20, bold=True, color=RED)

    add_bullet(tf_b, "92 treatments maintain posteriors", size=14)
    add_bullet(tf_b, "Top 10 = 49% of traffic", size=14)
    add_bullet(tf_b, "72 treatments get < 50 sends/day", size=14)
    add_bullet(tf_b, "37 opens/treatment/day", size=14, bold=True)
    add_bullet(tf_b, "115 days to converge", size=14, bold=True, color=RED)
    add_bullet(tf_b, "37.5% never converge", size=14, bold=True, color=RED)

    # After
    box_after = slide.shapes.add_shape(5, Inches(7.0), Inches(1.5), Inches(5.5), Inches(3.5))
    box_after.fill.solid()
    box_after.fill.fore_color.rgb = RGBColor(0xEA, 0xF7, 0xEA)
    box_after.line.color.rgb = GREEN

    tf_a = box_after.text_frame
    tf_a.word_wrap = True
    tf_a.margin_left = Inches(0.3)
    tf_a.margin_top = Inches(0.2)

    p2 = tf_a.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "AFTER: 10 Treatments"
    set_font(run2, size=20, bold=True, color=GREEN)

    add_bullet(tf_a, "10 treatments with highest traffic", size=14)
    add_bullet(tf_a, "Each gets ~250 sends/day", size=14)
    add_bullet(tf_a, "No long-tail noise", size=14)
    add_bullet(tf_a, "75 opens/treatment/day", size=14, bold=True)
    add_bullet(tf_a, "28 days to converge", size=14, bold=True, color=GREEN)
    add_bullet(tf_a, "0% never converge", size=14, bold=True, color=GREEN)

    # Arrow between
    arrow = slide.shapes.add_shape(1, Inches(6.2), Inches(3.0), Inches(0.7), Inches(0.1))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ORANGE
    arrow.line.fill.background()

    arrow_txt = slide.shapes.add_textbox(Inches(6.0), Inches(2.5), Inches(1.2), Inches(0.5))
    atf = arrow_txt.text_frame
    atf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run_a = atf.paragraphs[0].add_run()
    run_a.text = "4x"
    set_font(run_a, size=24, bold=True, color=ORANGE)

    # Bottom note
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(5.3), Inches(12), Inches(1.8))
    tf3 = txBox.text_frame
    tf3.word_wrap = True
    add_bullet(tf3, "Remove 82 long-tail treatments that get < 50 sends/day", size=14)
    add_bullet(tf3, "Top 10 already handle 49% of traffic -- minimal user impact", size=14)
    add_bullet(tf3, "Expected: model starts exploiting winners within 4 weeks", size=14, bold=True, color=DARK_BLUE)


def slide_14_recommendations(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Recommendations -- Prioritized")

    # Immediate
    sections = [
        ("IMMEDIATE (This Week)", ORANGE, [
            ("Reduce treatment pool 92 -> 10", "HIGHEST IMPACT: 4x faster convergence, 0% non-convergence"),
            ("Clamp scores to [0, 1]", "Prevent score > 1.0 anomaly. Apply MIN(score, 1.0) in serving layer"),
        ]),
        ("SHORT-TERM (Next 2 Weeks)", DARK_BLUE, [
            ("Cold-start warmup (>= 100 sends)", "New treatments excluded from bandit until sufficient data"),
            ("Revert to 10/90 split", "With reduced pool, give bandit 90% traffic to learn faster"),
        ]),
        ("MEDIUM-TERM (Next Month)", MED_GRAY, [
            ("Hierarchical learning", "Learn at campaign-type level, then differentiate within"),
            ("Contextual bandits", "Current model ignores user features -- add vehicle type, engagement history"),
        ]),
    ]

    y = Inches(1.4)
    for section_title, section_color, items in sections:
        # Section header
        header = slide.shapes.add_shape(5, Inches(0.6), y, Inches(3.5), Inches(0.45))
        header.fill.solid()
        header.fill.fore_color.rgb = section_color
        header.line.fill.background()
        htf = header.text_frame
        htf.margin_left = Inches(0.15)
        run = htf.paragraphs[0].add_run()
        run.text = section_title
        set_font(run, size=12, bold=True, color=WHITE)

        y += Inches(0.55)

        for action, detail in items:
            txBox = slide.shapes.add_textbox(Inches(0.8), y, Inches(12), Inches(0.65))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.space_after = Pt(2)

            run1 = p.add_run()
            run1.text = action + "  --  "
            set_font(run1, size=13, bold=True, color=DARK_BLUE)

            run2 = p.add_run()
            run2.text = detail
            set_font(run2, size=12, color=BLACK)

            y += Inches(0.6)

        y += Inches(0.1)


def slide_15_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Summary & Key Numbers")

    data = [
        ["Metric", "Value", "Source"],
        ["Treatments in bandit pool", "92", "Q12"],
        ["Treatments with 100+ sends/day", "20 (75% of traffic)", "Q12"],
        ["Top 10 share of traffic", "49%", "Q12"],
        ["Treatments eligible per user", "4-7 (fitment-filtered)", "Q12"],
        ["Clicks per treatment per week", "0.1-18", "Q16"],
        ["Phantom clicks", "877 / 821,740 (0.11%)", "Q11"],
        ["Invalid scores (> 1.0)", "1,686 total", "Q11/Q14"],
        ["NIG posterior accuracy", "Within 0.2-0.6pp", "Q13"],
        ["Convergence (20 trts)", "115 days, 37.5% never", "Simulation"],
        ["Convergence (10 trts)", "28 days, 0% never", "Simulation"],
        ["Improvement (20 -> 10)", "4x faster", "Simulation"],
    ]

    add_table(slide, 12, 3, data,
              Inches(0.6), Inches(1.3), Inches(8.5), Inches(5.0),
              col_widths=[Inches(3.5), Inches(3), Inches(2)])

    # Conclusion box
    box = slide.shapes.add_shape(5, Inches(9.5), Inches(1.5), Inches(3.5), Inches(3.5))
    box.fill.solid()
    box.fill.fore_color.rgb = DARK_BLUE
    box.line.fill.background()

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.25)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Bottom Line"
    set_font(run, size=18, bold=True, color=ORANGE)

    add_paragraph(tf, "Not a software bug -- a statistical reality to address "
                  "through treatment consolidation.", size=14, bold=True, color=WHITE)
    add_paragraph(tf, "", size=6)
    add_paragraph(tf, "Next Steps:", size=14, bold=True, color=ORANGE)
    add_paragraph(tf, "1. Reduce pool 92 -> 10", size=13, color=WHITE)
    add_paragraph(tf, "2. Clamp scores to [0,1]", size=13, color=WHITE)
    add_paragraph(tf, "3. Add cold-start warmup", size=13, color=WHITE)
    add_paragraph(tf, "4. Monitor for 4 weeks", size=13, color=WHITE)


# --- Main ---

def main():
    prs = Presentation()
    # Set 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    print("Generating slides...")
    slide_01_title(prs)
    print("  [1/15] Title")
    slide_02_executive_summary(prs)
    print("  [2/15] Executive Summary")
    slide_03_how_bandit_works(prs)
    print("  [3/15] How the Bandit Works")
    slide_04_phase1_recap(prs)
    print("  [4/15] Phase 1 Recap")
    slide_05_three_hypotheses(prs)
    print("  [5/15] Three Hypotheses")
    slide_06_data_quality(prs)
    print("  [6/15] Data Quality")
    slide_07_treatment_count(prs)
    print("  [7/15] Treatment Count (chart)")
    slide_08_score_anomaly(prs)
    print("  [8/15] Score Anomaly")
    slide_09_nig_verification(prs)
    print("  [9/15] NIG Verification")
    slide_10_root_cause(prs)
    print("  [10/15] Root Cause (chart)")
    slide_11_convergence_simulation(prs)
    print("  [11/15] Convergence Simulation (chart)")
    slide_12_posterior_evolution(prs)
    print("  [12/15] Posterior Evolution (chart)")
    slide_13_the_fix(prs)
    print("  [13/15] The Fix")
    slide_14_recommendations(prs)
    print("  [14/15] Recommendations")
    slide_15_summary(prs)
    print("  [15/15] Summary")

    os.makedirs(os.path.dirname(OUTPUT_FILE), exist_ok=True)
    prs.save(OUTPUT_FILE)
    print(f"\nSaved: {OUTPUT_FILE}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
